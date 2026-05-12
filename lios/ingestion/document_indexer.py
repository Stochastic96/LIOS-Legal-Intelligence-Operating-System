"""Document indexer — extract, chunk, and append uploaded documents to the retrieval corpus.

Supports PDF, DOCX, TXT, PPTX, and XLSX uploads. After indexing, the
HybridRetriever singleton is refreshed so new content is immediately searchable
without restarting the server.
"""

from __future__ import annotations

import json
from datetime import datetime, timezone
from hashlib import sha256
from pathlib import Path
from typing import Any

from lios.logging_setup import get_logger

logger = get_logger(__name__)

_CORPUS_PATH = Path("data/corpus/legal_chunks.jsonl")
_CHUNK_SIZE_WORDS = 400
_CHUNK_OVERLAP_WORDS = 50
SUPPORTED_UPLOAD_FORMATS: dict[str, str] = {
    "pdf": ".pdf",
    "docx": ".docx",
    "txt": ".txt",
    "pptx": ".pptx",
    "xlsx": ".xlsx",
}

_MIME_HINTS: dict[str, str] = {
    "pdf": "pdf",
    "docx": "wordprocessingml",
    "txt": "text/plain",
    "pptx": "presentationml",
    "xlsx": "spreadsheetml",
}


class UnsupportedDocumentFormatError(ValueError):
    """Raised when an upload format is unsupported."""


def supported_upload_extensions() -> list[str]:
    """Return supported upload extensions in API-facing order."""
    return list(SUPPORTED_UPLOAD_FORMATS.values())


def detect_upload_format(filename: str, content_type: str) -> str:
    """Resolve an upload format from file extension or content type."""
    fname_lower = (filename or "").lower()
    content_type_lower = (content_type or "").lower()

    for fmt, ext in SUPPORTED_UPLOAD_FORMATS.items():
        if fname_lower.endswith(ext):
            return fmt

    for fmt, hint in _MIME_HINTS.items():
        if hint in content_type_lower:
            return fmt

    supported = ", ".join(supported_upload_extensions())
    raise UnsupportedDocumentFormatError(
        f"Unsupported file format for '{filename or 'upload'}' "
        f"(content_type='{content_type or 'unknown'}'). "
        f"Supported formats: {supported}."
    )


def _extract_text_pdf(content: bytes) -> str:
    """Extract text from a PDF byte buffer using pypdf."""
    try:
        import io
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p.strip() for p in pages if p.strip())
    except ImportError:
        raise ImportError("pypdf is required for PDF uploads. Install with: pip install pypdf")
    except Exception as exc:
        raise ValueError(f"Failed to parse PDF: {exc}") from exc


def _extract_text_docx(content: bytes) -> str:
    """Extract text from a DOCX byte buffer using python-docx."""
    try:
        import io
        import docx  # type: ignore
        doc = docx.Document(io.BytesIO(content))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ImportError("python-docx is required for DOCX uploads. Install with: pip install python-docx")
    except Exception as exc:
        raise ValueError(f"Failed to parse DOCX: {exc}") from exc


def _extract_text_txt(content: bytes) -> str:
    """Extract text from a plain text byte buffer with explicit encoding fallbacks."""
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return content.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("Failed to decode TXT file using utf-8, utf-16, or latin-1.")


def _extract_text_pptx(content: bytes) -> str:
    """Extract text from a PPTX byte buffer using python-pptx."""
    try:
        import io
        from pptx import Presentation  # type: ignore

        presentation = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    lines.append(str(text).strip())
        return "\n\n".join(lines)
    except ImportError:
        raise ImportError("python-pptx is required for PPTX uploads. Install with: pip install python-pptx")
    except Exception as exc:
        raise ValueError(f"Failed to parse PPTX: {exc}") from exc


def _extract_text_xlsx(content: bytes) -> str:
    """Extract text from an XLSX byte buffer using openpyxl."""
    try:
        import io
        from openpyxl import load_workbook  # type: ignore

        workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
        sheet_blocks: list[str] = []
        for sheet in workbook.worksheets:
            row_lines: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if cells:
                    row_lines.append(" | ".join(cells))
            if row_lines:
                sheet_blocks.append(f"{sheet.title}\n" + "\n".join(row_lines))
        return "\n\n".join(sheet_blocks)
    except ImportError:
        raise ImportError("openpyxl is required for XLSX uploads. Install with: pip install openpyxl")
    except Exception as exc:
        raise ValueError(f"Failed to parse XLSX: {exc}") from exc


def _extract_text(content: bytes, filename: str, content_type: str) -> str:
    """Dispatch to the correct extractor based on validated upload format."""
    upload_format = detect_upload_format(filename, content_type)
    if upload_format == "pdf":
        return _extract_text_pdf(content)
    if upload_format == "docx":
        return _extract_text_docx(content)
    if upload_format == "txt":
        return _extract_text_txt(content)
    if upload_format == "pptx":
        return _extract_text_pptx(content)
    if upload_format == "xlsx":
        return _extract_text_xlsx(content)
    raise UnsupportedDocumentFormatError(f"No extractor registered for upload format: {upload_format}")


def _chunk_text(text: str, chunk_size: int = _CHUNK_SIZE_WORDS, overlap: int = _CHUNK_OVERLAP_WORDS) -> list[str]:
    """Split *text* into overlapping word-count chunks."""
    words = text.split()
    if not words:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(words):
        end = min(start + chunk_size, len(words))
        chunks.append(" ".join(words[start:end]))
        if end == len(words):
            break
        start += chunk_size - overlap
    return chunks


def _make_chunk_record(
    text: str,
    index: int,
    *,
    title: str,
    regulation: str,
    source_description: str,
    filename: str,
    doc_id: str,
) -> dict[str, Any]:
    """Build a legal_chunks.jsonl-compatible dict for one text chunk."""
    now = datetime.now(timezone.utc).isoformat()
    raw = f"{doc_id}|{index}|{text}".encode("utf-8")
    version_hash = sha256(raw).hexdigest()
    chunk_id = sha256(f"upload:{doc_id}:{index}".encode()).hexdigest()[:16]
    article_id = f"upload-{index + 1:03d}"

    return {
        "chunk_id": chunk_id,
        "source_url": "",
        "celex_or_doc_id": doc_id,
        "jurisdiction": "CUSTOM",
        "regulation": regulation,
        "article": article_id,
        "article_id": article_id,
        "published_date": now[:10],
        "effective_date": now[:10],
        "version_hash": version_hash,
        "ingestion_timestamp": now,
        "title": f"{title} — part {index + 1}",
        "text": text,
        "source": source_description or filename,
        "is_uploaded": True,
    }


def index_uploaded_document(
    content: bytes,
    filename: str,
    content_type: str,
    title: str,
    regulation: str,
    source_description: str,
) -> dict[str, Any]:
    """Extract, chunk, and append *content* to the retrieval corpus.

    Rebuilds the HybridRetriever's embedding index so uploaded content is
    immediately searchable without restarting the server.

    Returns:
        Dict with ``status``, ``chunks_added``, ``filename``.
    """
    text = _extract_text(content, filename, content_type)
    if not text.strip():
        return {"status": "error", "message": "No text could be extracted from the file.", "chunks_added": 0}

    doc_id = sha256(f"{filename}:{len(content)}".encode()).hexdigest()[:12]
    raw_chunks = _chunk_text(text)
    if not raw_chunks:
        return {"status": "error", "message": "Document is empty after chunking.", "chunks_added": 0}

    records = [
        _make_chunk_record(
            chunk, i,
            title=title,
            regulation=regulation,
            source_description=source_description,
            filename=filename,
            doc_id=doc_id,
        )
        for i, chunk in enumerate(raw_chunks)
    ]

    # Annotate each chunk with lawyer-lens metadata
    try:
        from lios.ingestion.lawyer_lens import annotate_chunk
        for record in records:
            annotate_chunk(record)
    except Exception as exc:
        logger.warning("Lawyer-lens annotation skipped: %s", exc)

    # Append to corpus JSONL
    _CORPUS_PATH.parent.mkdir(parents=True, exist_ok=True)
    with _CORPUS_PATH.open("a", encoding="utf-8") as f:
        for record in records:
            f.write(json.dumps(record, ensure_ascii=False) + "\n")

    logger.info("Indexed %d chunks from '%s' (regulation=%s)", len(records), filename, regulation)

    # Refresh the HybridRetriever singleton so new chunks are searchable.
    try:
        from lios.retrieval.hybrid_retriever import _retriever_singleton, _retriever_lock, HybridRetriever
        import lios.retrieval.hybrid_retriever as _hr_module
        with _retriever_lock:
            _hr_module._retriever_singleton = HybridRetriever()
        logger.info("HybridRetriever refreshed after upload.")
    except Exception as exc:
        logger.warning("Could not refresh retriever after upload: %s", exc)

    return {
        "status": "indexed",
        "filename": filename,
        "chunks_added": len(records),
        "regulation": regulation,
        "doc_id": doc_id,
    }
